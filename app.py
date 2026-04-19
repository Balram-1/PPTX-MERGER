import os
import uuid
from copy import deepcopy
from flask import Flask, request, render_template, send_file, redirect
from pptx import Presentation

UPLOAD_FOLDER = "uploads"
OUTPUT_FOLDER = "output"

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

app = Flask(__name__)
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER


def copy_slide(slide, target_prs):
    layout = target_prs.slide_layouts[6]  # blank
    new_slide = target_prs.slides.add_slide(layout)

    for shape in slide.shapes:
        new_el = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')


def merge_pptx(file_paths):
    merged = Presentation()

    # remove default slide
    if len(merged.slides) > 0:
        r_id = merged.slides._sldIdLst[0].rId
        merged.part.drop_rel(r_id)
        del merged.slides._sldIdLst[0]

    for path in file_paths:
        try:
            prs = Presentation(path)
            for slide in prs.slides:
                copy_slide(slide, merged)
        except Exception as e:
            print(f"Error processing {path}: {e}")

    output_file = os.path.join(OUTPUT_FOLDER, f"merged_{uuid.uuid4().hex}.pptx")
    merged.save(output_file)
    return output_file


@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        files = request.files.getlist("files")

        if not files or files[0].filename == "":
            return redirect("/")

        file_paths = []

        for f in files:
            if not f.filename.lower().endswith(".pptx"):
                continue

            unique_name = f"{uuid.uuid4().hex}_{f.filename}"
            path = os.path.join(UPLOAD_FOLDER, unique_name)
            f.save(path)
            file_paths.append(path)

        if not file_paths:
            return "No valid PPTX files uploaded"

        merged_file = merge_pptx(file_paths)

        # cleanup uploads after merge
        for path in file_paths:
            try:
                os.remove(path)
            except:
                pass

        return send_file(merged_file, as_attachment=True)

    return render_template("index.html")


if __name__ == "__main__":
    app.run(debug=True)